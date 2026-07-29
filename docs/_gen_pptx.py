# -*- coding: utf-8 -*-
"""
Generates ClientHub_Presentation.pptx — a clean 5-7 minute project-review deck.
Author of project: Suryakanta Pradhan.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x1F, 0x35, 0x5E)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
LIGHTBLUE = RGBColor(0xE8, 0xF0, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, text, size=18, color=NAVY, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def bullets(slide, x, y, w, h, items, size=18, color=GREY, gap=8, bullet_color=BLUE):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        bold_lead = None
        rest = item
        if isinstance(item, tuple):
            bold_lead, rest = item
        rb = p.add_run()
        rb.text = "•  "
        rb.font.size = Pt(size)
        rb.font.color.rgb = bullet_color
        rb.font.bold = True
        if bold_lead:
            r1 = p.add_run()
            r1.text = bold_lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
        r = p.add_run()
        r.text = rest
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def header(slide, title, num):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.07), ACCENT)
    txt(slide, Inches(0.5), 0, Inches(11.5), Inches(1.15), title,
        size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, Inches(12.3), 0, Inches(0.8), Inches(1.15), num,
        size=16, color=LIGHTBLUE, bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ============================================================
# 1. TITLE SLIDE
# ============================================================
s = add_slide()
bg(s, NAVY)
rect(s, 0, Inches(3.05), SW, Inches(0.06), ACCENT)
txt(s, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2), "ClientHub",
    size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.2), Inches(11.3), Inches(0.8),
    "AI-Powered Agency Client & Project Management Portal",
    size=22, color=LIGHTBLUE, italic=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.5),
    "Submitted by", size=16, color=GREY, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.05), Inches(11.3), Inches(0.7),
    "Suryakanta Pradhan", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.85), Inches(11.3), Inches(0.5),
    "Master of Computer Applications (MCA)  •  Full-Stack Web Project",
    size=15, color=LIGHTBLUE, align=PP_ALIGN.CENTER)

# ============================================================
# 2. PROBLEM STATEMENT
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "The Problem", "02")
bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5),
        [("Fragmented tools:", " agencies juggle spreadsheets, separate task boards, email, and billing apps."),
         ("No client transparency:", " clients must email to learn how their project is going."),
         ("Manual reporting:", " writing weekly status updates for every client is slow, repetitive work."),
         ("Disconnected billing:", " invoicing lives outside the project tool, detached from real work done."),
         ("Late risk detection:", " slipping deadlines are noticed only after they become a problem.")],
        size=20, gap=16)

# ============================================================
# 3. SOLUTION OVERVIEW
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "The Solution — ClientHub", "03")
txt(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.8),
    "One role-aware portal that unifies three products usually bought separately:",
    size=19, color=NAVY, bold=True)
cards = [("Services", "Clients, projects & tasks with a drag-and-drop Kanban workflow."),
         ("SaaS", "Multi-tenant, role-based dashboards for Admin, Team & Client."),
         ("Fintech", "Invoicing + online payments with secure webhook verification.")]
cw = Inches(3.7); gap = Inches(0.35); x0 = Inches(0.9); y0 = Inches(2.6)
for i, (t, d) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, y0, cw, Inches(2.6), LIGHTBLUE)
    rect(s, x, y0, cw, Inches(0.7), BLUE)
    txt(s, x, y0, cw, Inches(0.7), t, size=20, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.2), y0 + Inches(0.9), cw - Inches(0.4), Inches(1.6), d,
        size=16, color=GREY)
txt(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8),
    "+ Real-time notifications, analytics dashboards, and an AI weekly progress engine.",
    size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# 4-11. MODULE SLIDES
# ============================================================
module_slides = [
    ("Module 1 — Authentication & Access", "04",
     [("Secure login:", " email + password with bcrypt hashing."),
      ("JWT sessions:", " stateless tokens verified on every request."),
      ("Role-based access:", " admin / team_member / client guards."),
      ("Password reset:", " time-limited emailed link via Nodemailer."),
      ("Hardening:", " input validation + rate-limiting on auth routes.")]),
    ("Module 2 & 3 — Clients & Projects", "05",
     [("Client directory:", " contacts, notes, ImageKit logo upload."),
      ("Portal accounts:", " each client links to a login account."),
      ("Search & filter:", " by name, status, or project count."),
      ("Projects:", " budget, timeline, status under each client."),
      ("Auto progress:", " % computed from completed vs. total tasks.")]),
    ("Module 4 — Kanban Task Board", "06",
     [("Four columns:", " To Do / In Progress / Review / Done."),
      ("Drag & drop:", " move tasks between columns with @dnd-kit."),
      ("Collaboration:", " per-task comment threads."),
      ("Time tracking:", " log hours per task."),
      ("Activity log:", " who changed what, and when.")]),
    ("Module 5 — Invoicing & Payments", "07",
     [("Invoice builder:", " line items, auto tax & total."),
      ("Smart suggest:", " line items from logged hours."),
      ("PDF export:", " downloadable, professional invoices."),
      ("Pay Now:", " Razorpay test-mode checkout."),
      ("Secure webhook:", " HMAC-verified on the raw body, idempotent.")]),
    ("Module 6 — Real-Time Notifications", "08",
     [("Socket.io:", " live events end-to-end."),
      ("Triggers:", " task assigned, comment, payment, status change."),
      ("Bell + badge:", " unread count in the navbar."),
      ("Clickable list:", " jump straight to the related item."),
      ("Email fallback:", " for high-importance events.")]),
    ("Module 7 — Analytics Dashboards", "09",
     [("Admin view:", " revenue, pending payments, project counts."),
      ("Team workload:", " and top clients by revenue."),
      ("Client view:", " progress %, deadlines, invoice summary."),
      ("Aggregated in SQL:", " grouped in the database, not in JavaScript."),
      ("Recharts:", " clean charts from live data.")]),
]
for title, num, items in module_slides:
    s = add_slide(); bg(s, WHITE); header(s, title, num)
    bullets(s, Inches(0.8), Inches(1.75), Inches(11.7), Inches(5), items, size=20, gap=14)

# ============================================================
# 12. AI FEATURE HIGHLIGHT
# ============================================================
s = add_slide(); bg(s, NAVY);
rect(s, 0, 0, SW, Inches(1.15), RGBColor(0x12, 0x22, 0x40))
rect(s, 0, Inches(1.15), SW, Inches(0.07), ACCENT)
txt(s, Inches(0.5), 0, Inches(12), Inches(1.15), "★  AI Progress Summary — Core Differentiator",
    size=27, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
txt(s, Inches(12.3), 0, Inches(0.8), Inches(1.15), "10", size=16, color=LIGHTBLUE,
    bold=True, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)
flow = ["Aggregate weekly activity", "Build structured prompt", "Call Groq API",
        "Summary + risk flag", "Store & display"]
fw = Inches(2.3); fy = Inches(1.6); fx0 = Inches(0.55)
for i, step in enumerate(flow):
    x = fx0 + i * (fw + Inches(0.15))
    rect(s, x, fy, fw, Inches(1.05), LIGHTBLUE)
    txt(s, x + Inches(0.1), fy, fw - Inches(0.2), Inches(1.05), step,
        size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow) - 1:
        txt(s, x + fw - Inches(0.05), fy, Inches(0.25), Inches(1.05), "›",
            size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.8), Inches(3.1), Inches(11.7), Inches(2.6),
        [("Output:", " a friendly, client-ready summary of done / pending / blockers."),
         ("Risk flag:", " on-track  •  at-risk  •  delayed — auto-classified."),
         ("Automated:", " weekly node-cron job for every active project + on-demand button."),
         ("Resilient:", " template fallback, 'unknown' flag, and an 'AI unavailable' badge if the API fails.")],
        size=18, color=LIGHTBLUE, gap=12, bullet_color=ACCENT)

# ============================================================
# 13. TECH STACK
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "Technology Stack", "11")
col1 = [("Frontend:", " React, Tailwind, Recharts, @dnd-kit, Axios"),
        ("Backend:", " Node.js, Express"),
        ("Database:", " PostgreSQL + Prisma ORM"),
        ("Auth:", " JWT + bcrypt"),
        ("Real-time:", " Socket.io")]
col2 = [("AI:", " Groq API (Llama 3.3, OpenAI-compatible)"),
        ("Payments:", " Razorpay (test mode)"),
        ("Email / Media:", " Nodemailer / ImageKit"),
        ("Jobs / PDF:", " node-cron / pdfkit"),
        ("Quality:", " Vitest · ESLint · Prettier"),
        ("DevOps:", " Docker · GitHub Actions CI"),
        ("Deploy:", " Vercel · Render · Neon")]
bullets(s, Inches(0.8), Inches(1.8), Inches(5.9), Inches(5), col1, size=19, gap=18)
bullets(s, Inches(6.9), Inches(1.8), Inches(5.9), Inches(5), col2, size=19, gap=18)

# ============================================================
# 14. ARCHITECTURE
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "System Architecture", "12")
# three tiers
def tier(x, label, sub, color):
    rect(s, x, Inches(2.0), Inches(3.6), Inches(2.7), color)
    txt(s, x, Inches(2.1), Inches(3.6), Inches(0.7), label, size=20, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(2.9), Inches(3.2), Inches(1.7), sub, size=14,
        color=WHITE, align=PP_ALIGN.CENTER)
tier(Inches(0.7), "CLIENT", "React SPA\nTailwind · Recharts\ndnd-kit\nAxios + Socket.io", BLUE)
tier(Inches(4.85), "SERVER", "Express REST API\nRoutes → Middleware →\nControllers → Services\nSocket.io · Cron", NAVY)
tier(Inches(9.0), "DATA & SERVICES", "PostgreSQL (Prisma)\nGroq AI · Razorpay\nImageKit · SMTP", ACCENT)
txt(s, Inches(4.3), Inches(3.15), Inches(0.6), Inches(0.5), "→", size=30, color=GREY, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(8.45), Inches(3.15), Inches(0.6), Inches(0.5), "→", size=30, color=GREY, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1),
    "Decoupled 3-tier design. Secrets stay server-side; the browser never talks to external APIs directly. "
    "JWT secures REST; Socket.io rooms target live events to the right users.",
    size=16, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# ============================================================
# 15. SCREENSHOTS PLACEHOLDER
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "Live Demo / Screenshots", "13")
labels = ["Login + demo credentials", "Admin analytics dashboard", "Kanban board",
          "Invoice + PDF", "Client AI summary", "Notifications"]
gw = Inches(3.7); gh = Inches(1.95); gx0 = Inches(0.9); gy0 = Inches(1.7)
for i, lab in enumerate(labels):
    col = i % 3; row = i // 3
    x = gx0 + col * (gw + Inches(0.3))
    y = gy0 + row * (gh + Inches(0.4))
    rect(s, x, y, gw, gh, LIGHTBLUE)
    txt(s, x, y, gw, gh, "[ " + lab + " ]", size=15, color=BLUE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============================================================
# 15b. ENGINEERING QUALITY & DEVOPS
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "Security & Data Isolation", "14")
bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5),
        [("Passwords:", " bcrypt-hashed; plain text is never stored."),
         ("Tokens:", " signed JWTs with the algorithm pinned, re-validated against the "
          "database so a disabled account loses access at once."),
         ("Tenant isolation:", " a client requesting another client's project gets 403 — "
          "enforced on the server, so bypassing the UI does not help."),
         ("No privilege escalation:", " public signup hard-codes the team_member role."),
         ("Payments:", " an invoice is marked paid only after an HMAC signature check "
          "on the raw webhook body — never on the browser's word."),
         ("Defence in depth:", " Zod validation, Helmet headers, rate-limited login, "
          "2 MB image-only uploads, and no stack traces in production responses.")],
        size=19, gap=12)

# ============================================================
# 16b. ENGINEERING QUALITY
# ============================================================
s = add_slide(); bg(s, WHITE); header(s, "Engineering Quality & DevOps", "15")
bullets(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5),
        [("Automated tests:", " 44 Vitest tests — auth, calculations, the AI fallback, RBAC, "
          "the error handler and webhook verification — all passing."),
         ("Linting & formatting:", " ESLint (flat config) + Prettier — zero warnings on both apps."),
         ("Continuous Integration:", " GitHub Actions runs lint + tests + build on every push and PR."),
         ("Containerized:", " Docker + docker-compose bring up Postgres, API, and frontend in one command."),
         ("Performance:", " route-level code splitting took the initial bundle from ~748 KB "
          "to ~233 KB; search is debounced; responses and assets are gzipped."),
         ("Operability:", " env validation at boot, a health check that reports database "
          "reachability, structured logging, and graceful shutdown on SIGTERM."),
         ("Graceful failure:", " every external integration is optional and degrades instead "
          "of crashing — the AI falls back to a computed summary.")],
        size=18, gap=10)

# ============================================================
# 16. CLOSING
# ============================================================
s = add_slide(); bg(s, NAVY)
rect(s, 0, Inches(2.9), SW, Inches(0.06), ACCENT)
txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8), "Thank You",
    size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(3.2), Inches(10.9), Inches(1.6),
    "ClientHub — a full-stack, AI-powered SaaS demonstrating secure multi-role access, real-time "
    "collaboration, online payments, and production-grade AI reporting with graceful fallback.",
    size=18, color=LIGHTBLUE, italic=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6), "Suryakanta Pradhan",
    size=26, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.95), Inches(11.3), Inches(0.5),
    "MCA  •  Full-Stack Developer  •  Email: __________   GitHub: __________",
    size=15, color=LIGHTBLUE, align=PP_ALIGN.CENTER)

prs.save('clienthub/docs/ClientHub_Presentation.pptx')
print('PPTX saved OK — slides:', len(prs.slides._sldIdLst))
